from __future__ import annotations

import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable, Optional

from PIL import Image, ImageOps

LogCallback = Callable[[str], None]
ProgressCallback = Callable[[int, int, str], None]
CancelCallback = Callable[[], bool]


SUPPORTED_PRESENTATIONS = {".ppt", ".pptx", ".pptm", ".pps", ".ppsx"}
SUPPORTED_IMAGES = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}


def natural_sort_key(value: str | Path) -> list[object]:
    text = str(value)
    return [int(part) if part.isdigit() else part.casefold() for part in re.split(r"(\d+)", text)]


def safe_stem(name: str) -> str:
    value = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", name).strip().rstrip(".")
    return value or "未命名"


def find_soffice() -> Optional[str]:
    candidates = [
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    if sys.platform.startswith("win"):
        candidates.extend(
            [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            ]
        )
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return str(candidate)
    return None


def find_powerpoint_executable() -> Optional[str]:
    if not sys.platform.startswith("win"):
        return None
    try:
        import winreg

        locations = [
            (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\POWERPNT.EXE"),
            (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\App Paths\POWERPNT.EXE"),
            (winreg.HKEY_CURRENT_USER, r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\POWERPNT.EXE"),
        ]
        for hive, key_path in locations:
            try:
                with winreg.OpenKey(hive, key_path) as key:
                    value, _ = winreg.QueryValueEx(key, None)
                    if value and Path(value).exists():
                        return str(value)
            except OSError:
                continue
    except Exception:
        return None
    return None


def is_powerpoint_available() -> bool:
    if not sys.platform.startswith("win"):
        return False
    if find_powerpoint_executable():
        return True
    try:
        import winreg

        with winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, r"PowerPoint.Application\CLSID"):
            return True
    except Exception:
        return False


def detect_export_backends() -> dict[str, object]:
    powerpoint_path = find_powerpoint_executable()
    libreoffice_path = find_soffice()
    return {
        "powerpoint": is_powerpoint_available(),
        "powerpoint_path": powerpoint_path,
        "libreoffice": bool(libreoffice_path),
        "libreoffice_path": libreoffice_path,
    }


def _check_cancel(should_cancel: Optional[CancelCallback]) -> None:
    if should_cancel and should_cancel():
        raise InterruptedError("任务已取消")


def _save_pixmap_as_image(pix, destination: Path, image_format: str, jpg_quality: int) -> None:
    destination.parent.mkdir(parents=True, exist_ok=True)
    fmt = image_format.upper()
    if fmt == "PNG":
        pix.save(str(destination))
        return

    mode = "RGBA" if pix.alpha else "RGB"
    image = Image.frombytes(mode, (pix.width, pix.height), pix.samples)
    if mode == "RGBA":
        background = Image.new("RGB", image.size, "white")
        background.paste(image, mask=image.getchannel("A"))
        image = background
    image.save(destination, "JPEG", quality=max(1, min(100, jpg_quality)), optimize=True)


def _copy_or_convert_png(source: Path, destination: Path, image_format: str, jpg_quality: int) -> None:
    destination.parent.mkdir(parents=True, exist_ok=True)
    if image_format.upper() == "PNG":
        shutil.copy2(source, destination)
        return
    with Image.open(source) as opened:
        image = ImageOps.exif_transpose(opened)
        if image.mode in ("RGBA", "LA"):
            background = Image.new("RGB", image.size, "white")
            alpha = image.getchannel("A") if image.mode == "RGBA" else image.getchannel("A")
            background.paste(image.convert("RGB"), mask=alpha)
            image = background
        elif image.mode != "RGB":
            image = image.convert("RGB")
        image.save(destination, "JPEG", quality=max(1, min(100, jpg_quality)), optimize=True)


def _export_with_powerpoint(
    ppt_path: Path,
    temp_export_dir: Path,
    dpi: int,
    log: LogCallback,
    should_cancel: Optional[CancelCallback],
) -> list[Path]:
    if not sys.platform.startswith("win"):
        raise RuntimeError("PowerPoint 导出仅支持 Windows")

    try:
        import pythoncom
        import win32com.client
    except ImportError as exc:
        raise RuntimeError("缺少 pywin32，请先运行“安装依赖.bat”") from exc

    pythoncom.CoInitialize()
    app = None
    presentation = None
    try:
        _check_cancel(should_cancel)
        log("正在调用 Microsoft PowerPoint…")
        app = win32com.client.DispatchEx("PowerPoint.Application")
        try:
            app.DisplayAlerts = 0
            app.AutomationSecurity = 3
        except Exception:
            pass

        presentation = app.Presentations.Open(str(ppt_path.resolve()), True, False, False)
        slide_width_points = float(presentation.PageSetup.SlideWidth)
        slide_height_points = float(presentation.PageSetup.SlideHeight)
        pixel_width = max(1, round(slide_width_points / 72 * dpi))
        pixel_height = max(1, round(slide_height_points / 72 * dpi))

        temp_export_dir.mkdir(parents=True, exist_ok=True)
        presentation.Export(str(temp_export_dir.resolve()), "PNG", pixel_width, pixel_height)
        _check_cancel(should_cancel)

        exported = sorted(temp_export_dir.glob("*.png"), key=natural_sort_key)
        if not exported:
            exported = sorted(temp_export_dir.glob("*.PNG"), key=natural_sort_key)
        if not exported:
            raise RuntimeError("PowerPoint 未生成图片，请检查文件是否损坏或被保护")
        return exported
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


def _export_with_libreoffice(
    ppt_path: Path,
    temp_root: Path,
    dpi: int,
    image_format: str,
    jpg_quality: int,
    log: LogCallback,
    should_cancel: Optional[CancelCallback],
) -> list[Path]:
    soffice = find_soffice()
    if not soffice:
        raise RuntimeError("未检测到 LibreOffice")

    pdf_dir = temp_root / "pdf"
    profile_dir = temp_root / "lo_profile"
    render_dir = temp_root / "rendered"
    pdf_dir.mkdir(parents=True, exist_ok=True)
    profile_dir.mkdir(parents=True, exist_ok=True)
    render_dir.mkdir(parents=True, exist_ok=True)

    _check_cancel(should_cancel)
    log("正在调用 LibreOffice 转换为 PDF…")
    command = [
        soffice,
        "--headless",
        f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
        "--convert-to",
        "pdf",
        "--outdir",
        str(pdf_dir.resolve()),
        str(ppt_path.resolve()),
    ]
    creationflags = 0
    if sys.platform.startswith("win"):
        creationflags = getattr(subprocess, "CREATE_NO_WINDOW", 0)

    result = subprocess.run(
        command,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=240,
        creationflags=creationflags,
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice 转换失败：{result.stdout.strip()}")

    pdf_candidates = sorted(pdf_dir.glob("*.pdf"), key=natural_sort_key)
    if not pdf_candidates:
        raise RuntimeError(f"LibreOffice 未生成 PDF：{result.stdout.strip()}")
    pdf_path = pdf_candidates[0]

    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("缺少 PyMuPDF，请先运行“安装依赖.bat”") from exc

    log("正在渲染幻灯片图片…")
    document = fitz.open(pdf_path)
    try:
        pages: list[Path] = []
        matrix = fitz.Matrix(dpi / 72, dpi / 72)
        extension = ".png" if image_format.upper() == "PNG" else ".jpg"
        for index, page in enumerate(document, start=1):
            _check_cancel(should_cancel)
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            target = render_dir / f"page_{index:04d}{extension}"
            _save_pixmap_as_image(pix, target, image_format, jpg_quality)
            pages.append(target)
        return pages
    finally:
        document.close()


@dataclass(slots=True)
class PptExportOptions:
    output_dir: Path
    image_format: str = "PNG"
    dpi: int = 200
    jpg_quality: int = 92
    backend: str = "auto"  # auto / powerpoint / libreoffice
    separate_folder: bool = True


def export_presentations_to_images(
    presentation_paths: Iterable[str | Path],
    options: PptExportOptions,
    log: Optional[LogCallback] = None,
    progress: Optional[ProgressCallback] = None,
    should_cancel: Optional[CancelCallback] = None,
) -> list[Path]:
    logger = log or (lambda _message: None)
    files = [Path(path) for path in presentation_paths]
    if not files:
        raise ValueError("请先添加 PPT/PPTX 文件")
    for path in files:
        if not path.exists():
            raise FileNotFoundError(path)
        if path.suffix.lower() not in SUPPORTED_PRESENTATIONS:
            raise ValueError(f"不支持的演示文稿格式：{path.name}")

    image_format = options.image_format.upper()
    if image_format not in {"PNG", "JPG", "JPEG"}:
        raise ValueError("图片格式仅支持 PNG 或 JPG")
    if image_format == "JPEG":
        image_format = "JPG"
    dpi = max(72, min(600, int(options.dpi)))
    jpg_quality = max(1, min(100, int(options.jpg_quality)))
    options.output_dir.mkdir(parents=True, exist_ok=True)

    backends = detect_export_backends()
    requested = options.backend.lower()
    if requested not in {"auto", "powerpoint", "libreoffice"}:
        requested = "auto"

    all_outputs: list[Path] = []
    total_files = len(files)
    for file_index, ppt_path in enumerate(files, start=1):
        _check_cancel(should_cancel)
        logger(f"[{file_index}/{total_files}] 开始处理：{ppt_path.name}")
        if progress:
            progress(file_index - 1, total_files, ppt_path.name)

        stem = safe_stem(ppt_path.stem)
        destination_dir = options.output_dir / stem if options.separate_folder else options.output_dir
        destination_dir.mkdir(parents=True, exist_ok=True)
        extension = ".png" if image_format == "PNG" else ".jpg"

        with tempfile.TemporaryDirectory(prefix="ppt_image_converter_") as temp_name:
            temp_root = Path(temp_name)
            rendered: list[Path] = []
            used_backend = ""

            if requested in {"auto", "powerpoint"}:
                if backends["powerpoint"]:
                    try:
                        rendered = _export_with_powerpoint(
                            ppt_path,
                            temp_root / "powerpoint_export",
                            dpi,
                            logger,
                            should_cancel,
                        )
                        used_backend = "PowerPoint"
                    except Exception as exc:
                        if requested == "powerpoint":
                            raise
                        logger(f"PowerPoint 导出失败，尝试 LibreOffice：{exc}")
                elif requested == "powerpoint":
                    raise RuntimeError("未检测到 Microsoft PowerPoint")

            if not rendered and requested in {"auto", "libreoffice"}:
                if backends["libreoffice"]:
                    rendered = _export_with_libreoffice(
                        ppt_path,
                        temp_root,
                        dpi,
                        image_format,
                        jpg_quality,
                        logger,
                        should_cancel,
                    )
                    used_backend = "LibreOffice"
                elif requested == "libreoffice":
                    raise RuntimeError("未检测到 LibreOffice")

            if not rendered:
                raise RuntimeError("未找到可用导出引擎。请安装 PowerPoint 或 LibreOffice。")

            logger(f"使用 {used_backend}，共识别 {len(rendered)} 页")
            for slide_index, source_image in enumerate(rendered, start=1):
                _check_cancel(should_cancel)
                destination = destination_dir / f"{stem}_{slide_index:03d}{extension}"
                if used_backend == "PowerPoint":
                    _copy_or_convert_png(source_image, destination, image_format, jpg_quality)
                else:
                    if source_image.resolve() != destination.resolve():
                        shutil.copy2(source_image, destination)
                all_outputs.append(destination)

        logger(f"完成：{ppt_path.name} → {destination_dir}")
        if progress:
            progress(file_index, total_files, ppt_path.name)

    return all_outputs


@dataclass(slots=True)
class ImagesToPptOptions:
    output_path: Path
    slide_size: str = "16:9"  # 16:9 / 4:3 / A4横向 / A4竖向 / 按首图比例
    fit_mode: str = "contain"  # contain / cover
    background: str = "white"  # white / black


def _slide_size_inches(choice: str, first_image_size: tuple[int, int]) -> tuple[float, float]:
    presets = {
        "16:9": (13.333333, 7.5),
        "4:3": (10.0, 7.5),
        "A4横向": (11.6929, 8.2677),
        "A4竖向": (8.2677, 11.6929),
    }
    if choice in presets:
        return presets[choice]

    width_px, height_px = first_image_size
    if width_px <= 0 or height_px <= 0:
        return presets["16:9"]
    ratio = width_px / height_px
    if ratio >= 1:
        width_in = 13.333333
        height_in = width_in / ratio
    else:
        height_in = 13.333333
        width_in = height_in * ratio
    return max(1.0, width_in), max(1.0, height_in)


def _prepare_image_for_ppt(source: Path, temp_dir: Path, index: int) -> tuple[Path, int, int]:
    with Image.open(source) as opened:
        transposed = ImageOps.exif_transpose(opened)
        width, height = transposed.size
        needs_conversion = source.suffix.lower() not in {".png", ".jpg", ".jpeg"}
        has_exif_orientation = bool(opened.getexif().get(274, 1) != 1)
        if not needs_conversion and not has_exif_orientation:
            return source, width, height

        destination = temp_dir / f"normalized_{index:04d}.png"
        if transposed.mode == "P":
            transposed = transposed.convert("RGBA")
        elif transposed.mode not in {"RGB", "RGBA"}:
            transposed = transposed.convert("RGBA" if "A" in transposed.getbands() else "RGB")
        transposed.save(destination, "PNG")
        return destination, width, height


def _set_slide_background(slide, color_name: str) -> None:
    from pptx.dml.color import RGBColor

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0) if color_name == "black" else RGBColor(255, 255, 255)


def images_to_presentation(
    image_paths: Iterable[str | Path],
    options: ImagesToPptOptions,
    log: Optional[LogCallback] = None,
    progress: Optional[ProgressCallback] = None,
    should_cancel: Optional[CancelCallback] = None,
) -> Path:
    logger = log or (lambda _message: None)
    files = [Path(path) for path in image_paths]
    if not files:
        raise ValueError("请先添加图片")
    for path in files:
        if not path.exists():
            raise FileNotFoundError(path)
        if path.suffix.lower() not in SUPPORTED_IMAGES:
            raise ValueError(f"不支持的图片格式：{path.name}")

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx，请先运行“安装依赖.bat”") from exc

    output_path = options.output_path
    if output_path.suffix.lower() != ".pptx":
        output_path = output_path.with_suffix(".pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with Image.open(files[0]) as first_opened:
        first_image = ImageOps.exif_transpose(first_opened)
        first_size = first_image.size

    width_in, height_in = _slide_size_inches(options.slide_size, first_size)
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    blank_layout = prs.slide_layouts[6]

    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    fit_mode = options.fit_mode if options.fit_mode in {"contain", "cover"} else "contain"
    background = options.background if options.background in {"white", "black"} else "white"

    total = len(files)
    with tempfile.TemporaryDirectory(prefix="images_to_ppt_") as temp_name:
        temp_dir = Path(temp_name)
        for index, image_path in enumerate(files, start=1):
            _check_cancel(should_cancel)
            if progress:
                progress(index - 1, total, image_path.name)
            logger(f"[{index}/{total}] 添加图片：{image_path.name}")

            prepared_path, image_width, image_height = _prepare_image_for_ppt(image_path, temp_dir, index)
            slide = prs.slides.add_slide(blank_layout)
            _set_slide_background(slide, background)

            image_ratio = image_width / image_height
            slide_ratio = slide_width / slide_height

            if fit_mode == "contain":
                if image_ratio >= slide_ratio:
                    picture_width = slide_width
                    picture_height = int(slide_width / image_ratio)
                else:
                    picture_height = slide_height
                    picture_width = int(slide_height * image_ratio)
                left = int((slide_width - picture_width) / 2)
                top = int((slide_height - picture_height) / 2)
                slide.shapes.add_picture(str(prepared_path), left, top, width=picture_width, height=picture_height)
            else:
                picture = slide.shapes.add_picture(
                    str(prepared_path), 0, 0, width=slide_width, height=slide_height
                )
                if image_ratio > slide_ratio:
                    visible_fraction = slide_ratio / image_ratio
                    crop_each = max(0.0, (1.0 - visible_fraction) / 2.0)
                    picture.crop_left = crop_each
                    picture.crop_right = crop_each
                elif image_ratio < slide_ratio:
                    visible_fraction = image_ratio / slide_ratio
                    crop_each = max(0.0, (1.0 - visible_fraction) / 2.0)
                    picture.crop_top = crop_each
                    picture.crop_bottom = crop_each

            if progress:
                progress(index, total, image_path.name)

        _check_cancel(should_cancel)
        logger("正在保存 PPTX…")
        prs.save(output_path)

    logger(f"完成：{output_path}")
    return output_path
